# main.py
# Smilegate HRD | Mentoring Letter Auto Generator (Streamlit App)
# 실행: streamlit run main.py

import io
from datetime import date

import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

APP_TITLE = "멘토링 Letter 자동 생성기"
FIRST_SENTENCE_TEMPLATE = "{mentor} 멘토님, {mentee} 멘티의 멘토링 지원을 잘 부탁드립니다."

# --------- 폰트 사이즈 전역 설정 ---------
HEADER_FONT_SIZE = 14          # 맨 위 설명 문구
SECTION_TITLE_SIZE = 20        # "멘토에게", "활동 후기"
FIRST_SENTENCE_SIZE = 14       # test 멘토님, TEST 멘티...
BOX_TITLE_SIZE = 15            # "조직장 요청사항", "멘티 질문·고민", 우측 박스 제목
BODY_FONT_SIZE = 12            # 본문 텍스트
FOOTER_FONT_SIZE = 9           # 맨 아래 Mentor/Mentee/Date

# --------- 기본 문구 ---------
DEFAULT_REQUEST_TEXT = """1) 조직, 회사에 대한 이해
  - 조직의 방향성 및 구성에 대한 빠른 학습
  - 안정적으로 팀 문화에 적응할 수 있도록 도와주세요.
  - 업무적으로 편안하게 질문 할 수 있는 관계 형성이 되면 좋겠습니다.

2) 성장 및 업무 관련 지원
  - 팀 업무를 위해 사용 필요한 각종 시스템 및 프로세스에 대해 알려주세요.
  - 앞으로 맡아서 진행할 프로젝트 내 역할 분담"""

DEFAULT_MENTOR_NOTE = """▶ 리더 요청 사항 기반 활동한 내용을 간단하게 작성해주세요
▶ 추가적으로 조직장이 F/U이 필요한 사항을 작성해주세요.
   (ex 멘토링 활동간 멘티 궁금해 했으나, 답변을 못한 부분 or 요청한 사항)"""

THEME_COLOR = "#0B2B4C"  # 네이비 톤
RIGHT_BG = (237, 233, 226)
FONT_NAME = "Malgun Gothic"


def _add_textbox(
    slide,
    left_in,
    top_in,
    width_in,
    height_in,
    title,
    body,
    font_size_title=BOX_TITLE_SIZE,
    font_size_body=BODY_FONT_SIZE,
    bold_title=True,
):
    left = Inches(left_in)
    top = Inches(top_in)
    width = Inches(width_in)
    height = Inches(height_in)
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True

    # 제목
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(font_size_title)
    run.font.bold = bold_title
    run.font.name = FONT_NAME

    # 간격
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(2)

    # 본문
    for line in (body or "").splitlines():
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size_body)
        p.font.name = FONT_NAME
    return shape


def _add_rect(
    slide,
    left_in,
    top_in,
    width_in,
    height_in,
    fill_rgb=None,
    line_rgb=(180, 180, 180),
    line_width_pt=1.0,
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left_in),
        Inches(top_in),
        Inches(width_in),
        Inches(height_in),
    )
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = RGBColor(*line_rgb)
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape

def build_ppt(
    mentor,
    mentee,
    manager,
    first_sentence_template,
    request_text,
    use_default_request,
    qna_text,
    hide_qna_if_empty,
    mentor_note_text,
    logo_bytes,
    theme_color_hex,
):

    prs = Presentation()

    # 16:9 비율 고정
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # 전체 테두리 (그대로)
    _add_rect(slide, 0.4, 0.4, 12.5, 6.7, None, (80, 80, 80), 1.25)

    # 로고 (살짝 왼쪽)
    if logo_bytes:
        slide.shapes.add_picture(
            io.BytesIO(logo_bytes),
            Inches(0.55),
            Inches(0.55),
            height=Inches(0.45),
        )

    # 상단 설명 문구 → x: 1.0 -> 0.8
    header = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.55), Inches(11.2), Inches(0.5)
    )
    tf = header.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "멘토링 Letter는 멘토/멘티가 유의미한 멘토링이 되도록 참고할 수 있는 내용을 리더가 멘토에게 보내는 메시지 입니다."
    r.font.size = Pt(HEADER_FONT_SIZE)
    r.font.bold = True
    r.font.name = FONT_NAME

    # 섹션 제목 ("멘토에게", "활동 후기")
    # 둘 다 0.2inch 왼쪽으로 이동: 1.0 -> 0.8, 7.4 -> 7.2
    for (text, x) in [("멘토에게", 0.8), ("활동 후기", 7.2)]:
        box = slide.shapes.add_textbox(
            Inches(x), Inches(1.15), Inches(5.5), Inches(0.5)
        )
        tfb = box.text_frame
        tfb.clear()
        r = tfb.paragraphs[0].add_run()
        r.text = text
        r.font.size = Pt(SECTION_TITLE_SIZE)
        r.font.bold = True
        r.font.name = FONT_NAME

    # 첫 문장 → x: 1.0 -> 0.8
    sentence = first_sentence_template.format(
        mentor=mentor.strip(), mentee=mentee.strip()
    )
    box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.65), Inches(11.2), Inches(0.5)
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = sentence
    r.font.size = Pt(FIRST_SENTENCE_SIZE)
    r.font.name = FONT_NAME

    # 레이아웃 좌표 → left_x 1.0 -> 0.8 (전체 좌측으로 살짝 이동)
    left_x, top_y = 0.8, 2.2
    col_w, col_h = 6.0, 4.5

    # 우측 카드 배경도 함께 왼쪽으로: left_x + col_w + 0.25
    _add_rect(
        slide,
        left_x + col_w + 0.25,
        top_y,
        col_w,
        col_h,
        RIGHT_BG,
        (180, 180, 180),
        0.75,
    )

    # 좌측: 조직장 요청사항
    req = (request_text or "").strip()
    if use_default_request or len(req) < 5:
        req = DEFAULT_REQUEST_TEXT

    _add_textbox(
        slide,
        left_in=left_x,
        top_in=top_y,
        width_in=col_w,
        height_in=2.4,
        title="조직장 요청사항",
        body=req,
    )

    # 좌측: 멘티 질문·고민
    qna_text = qna_text or ""
    if not (hide_qna_if_empty and not qna_text.strip()):
        qna = qna_text.strip() or "(멘티 작성 예정)"
        _add_textbox(
            slide,
            left_in=left_x,
            top_in=top_y + 2.45,
            width_in=col_w,
            height_in=2.25,
            title="멘티 질문·고민",
            body=qna,
        )

    # 우측: 멘토 활동 후기 (같이 왼쪽으로 이동)
    _add_textbox(
        slide,
        left_in=left_x + col_w + 0.35,
        top_in=top_y + 0.15,
        width_in=col_w - 0.6,
        height_in=col_h - 0.3,
        title="멘토 활동 후기",
        body=mentor_note_text,
    )

    # 푸터 → 더 아래로: y 7.0 -> 7.25
    footer = slide.shapes.add_textbox(
        Inches(0.6), Inches(7.25), Inches(12.2), Inches(0.4)
    )
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"Mentor: {mentor}  |  Mentee: {mentee}  |  Date: {date.today():%Y.%m.%d}"
    r.font.size = Pt(FOOTER_FONT_SIZE)
    r.font.name = FONT_NAME

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio

def ui():
    st.set_page_config(page_title=APP_TITLE, page_icon="🧡", layout="wide")
    st.title(APP_TITLE)

    with st.sidebar:
        st.header("브랜딩 설정")
        theme = st.color_picker("포인트 색상", THEME_COLOR)
        logo_file = st.file_uploader(
            "로고 업로드 (PNG 권장)", type=["png", "jpg", "jpeg"]
        )
        st.caption("폰트는 시스템의 'Malgun Gothic'을 사용합니다.")

    col1, col2 = st.columns(2)
    with col1:
        st.subheader("인적 정보")
        mentor = st.text_input("멘토 이름")
        mentee = st.text_input("멘티 이름")
        manager = st.text_input("조직장(선택)")
        first_sentence_template = st.text_input(
            "첫 문장 템플릿", value=FIRST_SENTENCE_TEMPLATE
        )

        st.subheader("조직장 요청사항")
        request_text = st.text_area("요청사항 입력", height=200)
        use_default_request = st.checkbox(
            "비어있거나 짧으면 기본 양식 사용", value=True
        )

        st.subheader("멘티 질문·고민")
        qna_text = st.text_area("질문·고민 입력", height=140)
        hide_qna_if_empty = st.checkbox(
            "질문·고민이 없으면 해당 영역 삭제", value=True
        )

    with col2:
        st.subheader("멘토 활동 후기")
        mentor_note_text = st.text_area(
            "후기 가이드", value=DEFAULT_MENTOR_NOTE, height=260
        )

        if mentor and mentee:
            st.markdown(
                f"**미리보기:** {first_sentence_template.format(mentor=mentor, mentee=mentee)}"
            )
        else:
            st.caption("멘토/멘티 이름을 입력하면 첫 문장을 미리볼 수 있어요.")

    if st.button("PPT 생성 (다운로드)"):
        if not mentor or not mentee:
            st.error("멘토/멘티 이름은 필수입니다.")
            return
        logo_bytes = logo_file.read() if logo_file else None
        ppt_bytes = build_ppt(
            mentor,
            mentee,
            manager,
            first_sentence_template,
            request_text,
            use_default_request,
            qna_text,
            hide_qna_if_empty,
            mentor_note_text,
            logo_bytes,
            theme,
        )
        st.download_button(
            "PPT 다운로드",
            ppt_bytes,
            f"Mentoring_Letter_{mentee}_{mentor}.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )


if __name__ == "__main__":
    ui()

